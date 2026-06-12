"""
extractor.py
------------
Layer 3 execution: parse a quiz PDF or PPTX and return structured quiz data.
Handles variable formatting: numbered questions, lettered options, answer keys.
Self-heals: tries multiple parsing strategies before giving up.
"""

import re
import json
import logging
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Optional

log = logging.getLogger(__name__)


# ── Data models ──────────────────────────────────────────────────────────────

@dataclass
class Option:
    letter: str
    text: str

@dataclass
class Question:
    question: str
    options: list[Option]
    correct_answer: str = ""
    answer_explanation: str = ""

@dataclass
class QuizData:
    title: str
    category: str
    questions: list[Question]
    source_file: str
    status: str = "ok"          # ok | needs_review | failed
    notes: str = ""

    def to_dict(self):
        return asdict(self)


# ── Main entry point ──────────────────────────────────────────────────────────

def extract(file_path: str, default_category: str = "General Knowledge") -> QuizData:
    """
    Extract quiz data from a PDF or PPTX file.
    Returns a QuizData object. Never raises — sets status='failed' on error.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    try:
        if ext == ".pdf":
            raw_text = _extract_pdf(path)
        elif ext in (".pptx", ".ppt"):
            raw_text = _extract_pptx(path)
        else:
            return QuizData(
                title=path.stem, category=default_category,
                questions=[], source_file=str(path),
                status="failed", notes=f"Unsupported file type: {ext}"
            )
    except Exception as e:
        log.error(f"Extraction failed for {path.name}: {e}")
        return QuizData(
            title=path.stem, category=default_category,
            questions=[], source_file=str(path),
            status="failed", notes=str(e)
        )

    title = _find_title(raw_text, path.stem)
    category = _find_category(raw_text, default_category)
    questions = _parse_questions(raw_text)

    status = "ok"
    notes = ""
    if len(questions) < 3:
        status = "needs_review"
        notes = f"Only {len(questions)} question(s) found — may need manual check."
        log.warning(notes)

    return QuizData(
        title=title,
        category=category,
        questions=questions,
        source_file=str(path),
        status=status,
        notes=notes
    )


# ── PDF extraction ────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    """Try pdfplumber first (best for text PDFs), fall back to pypdf."""
    text = ""

    # Strategy 1: pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            parts = []
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        text = "\n".join(parts)
        if text.strip():
            log.info(f"pdfplumber extracted {len(text)} chars from {path.name}")
            return text
    except ImportError:
        log.warning("pdfplumber not installed — trying pypdf")
    except Exception as e:
        log.warning(f"pdfplumber failed: {e} — trying pypdf")

    # Strategy 2: pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = [page.extract_text() or "" for page in reader.pages]
        text = "\n".join(parts)
        if text.strip():
            log.info(f"pypdf extracted {len(text)} chars from {path.name}")
            return text
    except ImportError:
        log.warning("pypdf not installed — trying pdfminer")
    except Exception as e:
        log.warning(f"pypdf failed: {e}")

    # Strategy 3: pdfminer
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(str(path))
        log.info(f"pdfminer extracted {len(text)} chars from {path.name}")
        return text
    except ImportError:
        raise RuntimeError(
            "No PDF library found. Run: pip install pdfplumber pypdf pdfminer.six"
        )


# ── PPTX extraction ───────────────────────────────────────────────────────────

def _extract_pptx(path: Path) -> str:
    """Extract all text from all slides, preserving slide boundaries."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            slides_text.append(f"[SLIDE {i}]\n" + "\n".join(slide_lines))

    text = "\n\n".join(slides_text)
    log.info(f"pptx extracted {len(slides_text)} slides, {len(text)} chars from {path.name}")
    return text


# ── Text parsing helpers ──────────────────────────────────────────────────────

def _find_title(text: str, fallback: str) -> str:
    """Find the quiz title — usually the first non-empty line or a QUIZ/TITLE: label."""
    patterns = [
        r"(?i)(?:title|quiz name|quiz)[:\s]+(.+)",
        r"(?i)^(.{5,80})$",         # first short-ish line
    ]
    for pattern in patterns:
        m = re.search(pattern, text, re.MULTILINE)
        if m:
            candidate = m.group(1).strip().rstrip("?!.")
            if 4 < len(candidate) < 100:
                return candidate
    return fallback.replace("_", " ").replace("-", " ").title()


def _find_category(text: str, fallback: str) -> str:
    """Detect topic category from known keywords in the text."""
    categories = {
        "Science":          r"\b(science|biology|chemistry|physics|anatomy|periodic)\b",
        "History":          r"\b(history|historical|century|war|empire|dynasty|ancient)\b",
        "Geography":        r"\b(geography|capital|country|continent|river|mountain|ocean)\b",
        "Sports":           r"\b(sport|football|cricket|tennis|olympic|athlete|stadium)\b",
        "Technology":       r"\b(technology|computer|software|internet|coding|programming|AI)\b",
        "Entertainment":    r"\b(movie|film|music|actor|singer|celebrity|oscar|grammy)\b",
        "Mathematics":      r"\b(math|algebra|geometry|calculus|equation|fraction)\b",
    }
    text_lower = text.lower()
    scores = {}
    for cat, pattern in categories.items():
        matches = len(re.findall(pattern, text_lower))
        if matches:
            scores[cat] = matches
    if scores:
        return max(scores, key=scores.get)
    return fallback


def _parse_questions(text: str) -> list[Question]:
    """
    Parse questions and answer options from raw text.
    Tries multiple patterns to handle different quiz formats.
    """
    questions = []

    # Pattern 1: Numbered questions with lettered options
    # Q1. What is... / 1. What is... / Question 1: ...
    q_pattern = re.compile(
        r"(?:Q\.?\s*|Question\s*)?"
        r"(\d+)[.)]\s+"
        r"(.+?)"                              # question text
        r"(?=\n\s*[A-Da-d][.)]\s|\Z)",
        re.DOTALL
    )
    opt_pattern = re.compile(r"([A-Da-d])[.)]\s+(.+?)(?=\n\s*[A-Da-d][.)]|\n\s*\d+[.)]|\Z)", re.DOTALL)
    answer_pattern = re.compile(r"(?i)(?:ans(?:wer)?|correct)[:\s]+([A-D])\b")

    blocks = re.split(r"\n(?=(?:Q\.?\s*|Question\s*)?\d+[.)]\s)", text)

    for block in blocks:
        block = block.strip()
        if not block:
            continue

        q_match = re.match(r"(?:Q\.?\s*|Question\s*)?(\d+)[.)]\s+(.+)", block, re.DOTALL)
        if not q_match:
            continue

        q_text_raw = q_match.group(2)

        # Split question text from options
        opt_split = re.split(r"\n\s*[A-Da-d][.)]\s", q_text_raw, maxsplit=1)
        q_text = opt_split[0].strip().replace("\n", " ")

        # Extract options
        opts_raw = block[q_match.start(2):]
        options = []
        for om in opt_pattern.finditer(opts_raw):
            options.append(Option(
                letter=om.group(1).upper(),
                text=om.group(2).strip().replace("\n", " ")
            ))

        if not q_text or len(q_text) < 5:
            continue

        # Find answer
        ans_match = answer_pattern.search(block)
        correct = ans_match.group(1).upper() if ans_match else ""

        questions.append(Question(
            question=q_text,
            options=options,
            correct_answer=correct
        ))

    # Pattern 2: Slide-per-question format (PPTX)
    if not questions:
        questions = _parse_slide_format(text)

    log.info(f"Parsed {len(questions)} questions")
    return questions


def _parse_slide_format(text: str) -> list[Question]:
    """Fallback parser for PPTX where each slide is one question."""
    questions = []
    slides = re.split(r"\[SLIDE \d+\]", text)

    for slide in slides:
        lines = [l.strip() for l in slide.strip().split("\n") if l.strip()]
        if len(lines) < 2:
            continue

        q_text = lines[0]
        options = []
        correct = ""

        for line in lines[1:]:
            m = re.match(r"([A-Da-d])[.)]\s*(.+)", line)
            if m:
                options.append(Option(letter=m.group(1).upper(), text=m.group(2)))
            ans_m = re.match(r"(?i)(?:answer|ans)[:\s]+([A-D])", line)
            if ans_m:
                correct = ans_m.group(1).upper()

        if q_text and len(q_text) > 5:
            questions.append(Question(
                question=q_text,
                options=options,
                correct_answer=correct
            ))

    return questions


# ── CLI test ──────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python extractor.py <quiz_file.pdf|.pptx>")
        sys.exit(1)

    logging.basicConfig(level=logging.INFO)
    result = extract(sys.argv[1])
    print(json.dumps(result.to_dict(), indent=2))
