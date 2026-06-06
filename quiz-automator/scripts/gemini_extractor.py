"""
gemini_extractor.py — Vision-AI extractor using Google Gemini.

Reads slide images directly with Gemini's vision model. Far more accurate
than traditional OCR because the model UNDERSTANDS the slide:
  - Knows what's a question vs answer vs decoration
  - Reads stylized text and decorative backgrounds correctly
  - Handles "SOMEONE FAMOUS" picture-ID slides
  - Adapts to any quiz format automatically (no per-format coding)

Used by extractor.py when:
  - config.json has gemini_api_key set
  - The file has no extractable text (image-based PPTX/PDF)
"""
import os
import io
import json
import sys
import time
import base64
import tempfile
import shutil
from pathlib import Path

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR   = os.path.dirname(SCRIPT_DIR)


# ── Gemini setup ─────────────────────────────────────────────────────────────

_model = None

def get_model(api_key):
    global _model
    if _model is None:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        # Use Gemini 2.0 Flash — fast, accurate, free tier supports ~15 RPM, 1500/day
        _model = genai.GenerativeModel("gemini-2.0-flash")
    return _model


# ── Slide rendering ──────────────────────────────────────────────────────────

def render_pptx_slides(pptx_path, temp_dir):
    """
    Extract the largest image from each PPTX slide.
    Returns list of (slide_image_path, slide_index) tuples.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image

    prs = Presentation(pptx_path)
    paths = []
    for idx, slide in enumerate(prs.slides):
        images = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not images:
            paths.append((None, idx))
            continue
        biggest = max(images, key=lambda s: s.width * s.height)
        img = Image.open(io.BytesIO(biggest.image.blob))
        out = os.path.join(temp_dir, f"slide_{idx:04d}.png")
        # Convert to RGB and save (Gemini wants standard PNG/JPEG)
        if img.mode != "RGB":
            img = img.convert("RGB")
        img.save(out, "PNG")
        paths.append((out, idx))
    return paths


def render_pdf_slides(pdf_path, temp_dir, dpi=150):
    """Render each PDF page as PNG."""
    import fitz
    doc = fitz.open(pdf_path)
    paths = []
    for i, page in enumerate(doc):
        out = os.path.join(temp_dir, f"slide_{i:04d}.png")
        page.get_pixmap(dpi=dpi).save(out)
        paths.append((out, i))
    doc.close()
    return paths


# ── Gemini extraction ────────────────────────────────────────────────────────

# Single prompt that asks Gemini to look at the whole quiz at once and pull out
# Q&A pairs. We send slides in batches to stay within rate limits.

CLASSIFY_PROMPT = """You are processing a quiz presentation slide-by-slide.

For each image I send you, classify it as ONE of:
- QUESTION: a slide containing a quiz question (text or image-based)
- ANSWER:  a slide revealing the answer to a question
- TITLE:   title slide / round header / instructions / acknowledgements
- BLANK:   empty or unreadable

Respond with ONLY the single classification word (QUESTION / ANSWER / TITLE / BLANK).
"""

EXTRACT_PROMPT = """Look at these two consecutive slides from a quiz.
Slide 1 is a QUESTION slide. Slide 2 is the ANSWER slide.

Extract the question and answer in JSON format:
{
  "question": "the full question text, cleaned up — no slide numbers, no decorative noise",
  "answer":   "the answer (just the answer, not the explanation)",
  "funda":    "any extra explanation/trivia text on the answer slide, or empty string",
  "is_visual": true if the question slide is a photo of a person/object asking 'who/what is this', otherwise false
}

If the question slide just says 'answer on next slide' or shows a photo (no text question), set "question" to "Identify this." and is_visual to true.

Respond with ONLY the JSON object, no markdown fences.
"""


def classify_slide(model, image_path):
    """Ask Gemini to classify a single slide. Returns one of: QUESTION/ANSWER/TITLE/BLANK."""
    from PIL import Image
    img = Image.open(image_path)
    try:
        response = model.generate_content([CLASSIFY_PROMPT, img])
        label = response.text.strip().upper()
        for valid in ("QUESTION", "ANSWER", "TITLE", "BLANK"):
            if valid in label:
                return valid
        return "BLANK"
    except Exception as e:
        print(f"  Classify error on {os.path.basename(image_path)}: {e}", file=sys.stderr, flush=True)
        return "BLANK"


def extract_qa_pair(model, q_image_path, a_image_path):
    """Extract Q+A from a pair of slide images. Returns dict or None."""
    from PIL import Image
    q_img = Image.open(q_image_path)
    a_img = Image.open(a_image_path)
    try:
        response = model.generate_content([EXTRACT_PROMPT, q_img, a_img])
        text = response.text.strip()
        # Strip markdown fences if present
        if text.startswith("```"):
            text = text.split("```", 2)[1]
            if text.startswith("json"):
                text = text[4:]
            text = text.rsplit("```", 1)[0]
        data = json.loads(text.strip())
        return data
    except Exception as e:
        print(f"  Extract error: {e}", file=sys.stderr, flush=True)
        return None


# ── topic detection (reuse from ocr_extractor) ───────────────────────────────

def detect_topic(q_text, a_text, topic_keywords):
    combined = (q_text + " " + a_text).lower()
    scores = {
        topic: sum(1 for kw in kws if kw.lower() in combined)
        for topic, kws in topic_keywords.items()
    }
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "general"


# ── Main pipeline ────────────────────────────────────────────────────────────

def extract_from_image_file(path, config, show_progress=True):
    """
    Full Gemini-based extraction pipeline.
    Returns list of question dicts in extractor.py's format.
    """
    api_key = config.get("gemini_api_key", "").strip()
    if not api_key:
        raise RuntimeError("gemini_api_key not set in config.json")

    ext            = os.path.splitext(path)[1].lower()
    topic_keywords = config.get("topic_keywords", {})
    temp_dir       = tempfile.mkdtemp(prefix="quiz_gemini_")

    # Step 1 — render slides
    if show_progress:
        print("Rendering slides...", file=sys.stderr, flush=True)

    if ext == ".pdf":
        slides = render_pdf_slides(path, temp_dir)
    elif ext in (".pptx", ".ppt"):
        slides = render_pptx_slides(path, temp_dir)
    else:
        raise RuntimeError(f"Unsupported format: {ext}")

    total = len(slides)
    if show_progress:
        print(f"Got {total} slides. Connecting to Gemini...", file=sys.stderr, flush=True)

    model = get_model(api_key)

    # Step 2 — classify each slide
    if show_progress:
        print(f"Classifying slides via Gemini...", file=sys.stderr, flush=True)

    classifications = []
    for i, (img_path, idx) in enumerate(slides):
        if img_path is None:
            classifications.append("BLANK")
            continue
        label = classify_slide(model, img_path)
        classifications.append(label)
        if show_progress and (i % 5 == 0 or i == total - 1):
            print(f"  {i+1}/{total} ({label})", file=sys.stderr, flush=True)
        # Throttle: free tier ~15 RPM. Wait 4.2 seconds between calls.
        time.sleep(4.2)

    # Step 3 — pair Q→A and extract
    questions = []
    pairs = []
    i = 0
    while i < len(classifications):
        if classifications[i] == "QUESTION":
            # find next ANSWER
            j = i + 1
            while j < len(classifications) and classifications[j] not in ("ANSWER", "QUESTION"):
                j += 1
            if j < len(classifications) and classifications[j] == "ANSWER":
                pairs.append((i, j))
                i = j + 1
            else:
                i += 1
        else:
            i += 1

    if show_progress:
        print(f"Found {len(pairs)} Q→A pairs. Extracting content...", file=sys.stderr, flush=True)

    for k, (qi, ai) in enumerate(pairs):
        q_path = slides[qi][0]
        a_path = slides[ai][0]
        if not q_path or not a_path:
            continue

        if show_progress:
            print(f"  Q{k+1}/{len(pairs)} (slides {qi+1}→{ai+1})...", file=sys.stderr, flush=True)

        data = extract_qa_pair(model, q_path, a_path)
        time.sleep(4.2)
        if not data:
            continue

        q_text = (data.get("question") or "").strip()
        a_text = (data.get("answer") or "").strip()
        f_text = (data.get("funda") or "").strip()
        is_vis = bool(data.get("is_visual"))

        if not q_text or not a_text:
            continue

        topic = detect_topic(q_text, a_text, topic_keywords)
        questions.append({
            "question_text":         q_text,
            "answer_text":           a_text,
            "funda_text":            f_text,
            "topic":                 topic,
            "has_image_placeholder": is_vis,
            "question_image_path":   q_path if is_vis else None,
            "answer_image_path":     a_path,
        })

    if show_progress:
        print(f"Done. {len(questions)} questions extracted via Gemini.", file=sys.stderr, flush=True)

    # temp_dir kept — questions_updater needs the images
    return questions


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python gemini_extractor.py <pptx_or_pdf>", file=sys.stderr)
        sys.exit(1)

    with open(os.path.join(ROOT_DIR, "config.json"), "r", encoding="utf-8") as f:
        cfg = json.load(f)

    qs = extract_from_image_file(sys.argv[1], cfg)
    print(json.dumps({
        "title":     os.path.splitext(os.path.basename(sys.argv[1]))[0],
        "questions": qs,
    }, ensure_ascii=False, indent=2))
