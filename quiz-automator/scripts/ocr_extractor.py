"""
ocr_extractor.py — OCR-based extractor for PDFs/PPTXs where every slide is an image.
Called automatically by extractor.py when normal text extraction yields nothing.

Improvements over v1:
- Region-based OCR: answer slides are scanned in two strips (top-title + bottom-caption)
  to avoid reading noise text embedded inside images (posters, poems, etc.)
- Smart answer line scoring: short, clean English phrases are preferred
- Answer image extraction: the image on each answer slide is saved and linked
"""
import os
import re
import sys
import json
import tempfile
import shutil
import unicodedata

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR   = os.path.dirname(SCRIPT_DIR)


# ── OCR engine ───────────────────────────────────────────────────────────────

_reader = None  # cached so model loads only once per process

def get_reader():
    global _reader
    if _reader is None:
        import easyocr
        _reader = easyocr.Reader(["en"], verbose=False)
    return _reader


def ocr_pil_image(pil_img):
    """Run OCR on a PIL Image object. Returns text string."""
    import numpy as np
    reader = get_reader()
    arr = np.array(pil_img.convert("RGB"))
    results = reader.readtext(arr, detail=0, paragraph=True)
    return "\n".join(results)


def ocr_image(image_path):
    """Run OCR on a file path. Returns text string."""
    from PIL import Image
    img = Image.open(image_path)
    return ocr_pil_image(img)


# ── slide rendering ──────────────────────────────────────────────────────────

def render_pdf_slides(pdf_path, temp_dir, dpi=180):
    """Render each PDF page as a PNG. Returns list of (image_path, None) tuples.
    Second element reserved for embedded image path (not applicable for PDF)."""
    import fitz
    doc = fitz.open(pdf_path)
    paths = []
    for i, page in enumerate(doc):
        out = os.path.join(temp_dir, f"slide_{i:04d}.png")
        page.get_pixmap(dpi=dpi).save(out)
        paths.append((out, None))
    doc.close()
    return paths


def render_pptx_slides(pptx_path, temp_dir, dpi=180):
    """
    For each slide in the PPTX:
    - Extract the largest picture shape as the "slide render"
    - Also save that picture separately as the "embedded image"
    Returns list of (slide_render_path, embedded_image_path) tuples.
    embedded_image_path is the raw image from the slide (used as answer image).
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io

    prs   = Presentation(pptx_path)
    paths = []

    for slide_idx, slide in enumerate(prs.slides):
        best_img_bytes = None
        best_area      = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                area = shape.width * shape.height
                if area > best_area:
                    best_area      = area
                    best_img_bytes = shape.image.blob

        if best_img_bytes:
            # Slide render (used for OCR)
            slide_out = os.path.join(temp_dir, f"slide_{slide_idx:04d}.png")
            img = Image.open(io.BytesIO(best_img_bytes))
            img.save(slide_out, "PNG")

            # Embedded image (saved separately for use as answer image)
            emb_out = os.path.join(temp_dir, f"embed_{slide_idx:04d}.png")
            img.save(emb_out, "PNG")

            paths.append((slide_out, emb_out))
        else:
            paths.append((None, None))

    return paths


# ── region-based OCR helpers ─────────────────────────────────────────────────

def ocr_region(image_path, top_frac=0.0, bottom_frac=1.0, left_frac=0.0, right_frac=1.0):
    """
    OCR only a rectangular region of an image.
    Fractions are 0.0–1.0 of image dimensions.
    """
    from PIL import Image
    img = Image.open(image_path)
    w, h = img.size
    box = (
        int(left_frac  * w),
        int(top_frac   * h),
        int(right_frac * w),
        int(bottom_frac * h),
    )
    cropped = img.crop(box)
    return ocr_pil_image(cropped)


def ocr_answer_slide(image_path):
    """
    Smart OCR for answer slides.
    Scans two regions and picks the cleanest short result:
      - Title strip: just below the ANSWER banner (~15%–38% of height)
      - Caption strip: bottom of slide (~72%–100% of height)
    Falls back to full-slide OCR if both regions are empty.
    """
    # Title strip (answer name right after the ANSWER banner)
    title_text   = ocr_region(image_path, top_frac=0.13, bottom_frac=0.40)
    # Caption strip (label at the bottom)
    caption_text = ocr_region(image_path, top_frac=0.72, bottom_frac=1.00)

    title_clean   = smart_answer_line(title_text)
    caption_clean = smart_answer_line(caption_text)

    # Prefer whichever is shorter and looks cleaner
    def quality(t):
        words = t.split()
        if not words:
            return 0
        ascii_ratio = sum(1 for c in t if c.isascii() and (c.isalpha() or c.isspace())) / len(t)
        length_bonus = max(0, 8 - len(words))  # shorter is usually better
        return ascii_ratio * 5 + length_bonus

    if quality(title_clean) >= quality(caption_clean) and title_clean:
        best = title_clean
    elif caption_clean:
        best = caption_clean
    else:
        # Fall back to full slide
        full_text = ocr_image(image_path)
        best = smart_answer_line(full_text)

    return best


def ocr_question_slide(image_path):
    """
    OCR for question slides.
    Scan the full width — cropping was cutting off the first letter of every line.
    clean_question_text() already strips the Q. N label from the text.
    """
    text = ocr_region(image_path, left_frac=0.0)
    return clean_question_text(text)


# ── slide classification ──────────────────────────────────────────────────────

def is_answer_slide(text):
    """True if slide OCR starts with ANSWER."""
    for line in text.strip().splitlines()[:4]:
        if re.match(r"^ANSWER\s*$", line.strip(), re.IGNORECASE):
            return True
    return False


def is_junk_slide(q_text):
    """True if slide is clearly not a real question (section header, instructions, etc.)."""
    junk_patterns = [
        r"answer can be found",
        r"can be found on the next slide",
        r"someone famous",
        r"^prelims?\s*$",
        r"^finals?\s*$",
        r"^round\s*\d",
        r"^manifesto",
        r"^section\s*\d",
        r"^\d+\s+questions?,\s*\d+\s+points?",
        r"no negative marking",
        r"hammered question",
        r"^thank",
        r"^quiz show\s*$",
        r"^visual connect",
    ]
    lowered = q_text.lower()
    return any(re.search(p, lowered) for p in junk_patterns)


def is_visual_question(q_text):
    """
    True if the question slide is a visual-identification question —
    e.g. "SOMEONE FAMOUS", "answer on next slide".
    These have no readable question text but the answer slide has the name.
    We keep them with a generic question text instead of discarding them.
    """
    visual_patterns = [
        r"someone famous",
        r"can be found on the next slide",
        r"answer can be found",
        r"visual connect",
        r"identify (this|the)",
        r"^who (is|am|are) (this|i|these)",
    ]
    lowered = q_text.lower().strip()
    return any(re.search(p, lowered) for p in visual_patterns)


# ── text cleaning ─────────────────────────────────────────────────────────────

def ascii_ratio(text):
    if not text:
        return 0
    return sum(1 for c in text if c.isascii() and (c.isalpha() or c.isspace())) / len(text)


def smart_answer_line(text):
    """
    From raw OCR text of an answer region, extract the best single answer phrase.
    Strategy:
      1. Remove ANSWER header lines
      2. Score each line: prefer short, high-ASCII, English-looking lines
      3. Return best line (up to 120 chars)
    """
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    lines = [l for l in lines if not re.match(r"^ANSWER\s*$", l, re.IGNORECASE)]

    if not lines:
        return ""

    def score(line):
        words = line.split()
        if not words:
            return -1
        ar   = ascii_ratio(line)
        wlen = len(words)
        # Heavily penalise very long lines (poem content, etc.)
        length_penalty = max(0, wlen - 8) * 0.5
        return ar * 6 - length_penalty

    # Filter lines with at least 30% ASCII (remove pure Cyrillic/noise)
    clean_lines = [l for l in lines if ascii_ratio(l) >= 0.30]
    if not clean_lines:
        clean_lines = lines  # fallback

    best = max(clean_lines, key=score)
    # Normalise
    best = "".join(c for c in best if unicodedata.category(c) != "Cf")
    return " ".join(best.split())[:200]


def clean_question_text(text):
    """Strip Q. N labels and normalize question text."""
    lines = text.strip().splitlines()
    cleaned = []
    for line in lines:
        s = line.strip()
        # Skip standalone Q.N label lines  e.g. "Q. 3" or "Q3"
        if re.match(r"^Q[\.\s]*\d*\s*$", s, re.IGNORECASE):
            continue
        # Skip lines that are ONLY a number (slide/question number artefacts)
        if re.match(r"^\d{1,2}\.?\s*$", s):
            continue
        # Remove inline Q.N prefix at start of line e.g. "Q. 3  What is..."
        s = re.sub(r"^Q[\.\s]*\d+\s+", "", s)
        if s:
            cleaned.append(s)
    result = " ".join(cleaned)
    result = "".join(c for c in result if unicodedata.category(c) != "Cf")
    return " ".join(result.split())


# ── topic detection ──────────────────────────────────────────────────────────

def detect_topic(q_text, a_text, topic_keywords):
    combined = (q_text + " " + a_text).lower()
    scores = {
        topic: sum(1 for kw in kws if kw.lower() in combined)
        for topic, kws in topic_keywords.items()
    }
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "general"


# ── main pipeline ────────────────────────────────────────────────────────────

def extract_from_image_file(path, config, show_progress=True):
    """
    Full OCR pipeline for an image-slide PDF or PPTX.
    Returns list of question dicts in extractor.py's output format.
    """
    ext            = os.path.splitext(path)[1].lower()
    topic_keywords = config.get("topic_keywords", {})
    temp_dir       = tempfile.mkdtemp(prefix="quiz_ocr_")

    # ── Step 1: render slides ────────────────────────────────────────────────
    if show_progress:
        print("Rendering slides...", file=sys.stderr, flush=True)

    if ext == ".pdf":
        slide_tuples = render_pdf_slides(path, temp_dir)
    elif ext in (".pptx", ".ppt"):
        slide_tuples = render_pptx_slides(path, temp_dir)
    else:
        raise RuntimeError(f"Unsupported format for OCR extraction: {ext}")

    total = len(slide_tuples)
    if show_progress:
        print(f"OCR scanning {total} slides (this takes a few minutes)...", file=sys.stderr)

    # ── Step 2: fast classification pass — only scan top strip of each slide ──
    # We only need to detect the word "ANSWER" at the top to classify slides.
    # Scanning just the top 20% is ~5x faster than full-slide OCR.
    if show_progress:
        print(f"Pass 1 (fast): classifying {total} slides...", file=sys.stderr, flush=True)

    answer_indices = set()
    for idx, (slide_path, _emb) in enumerate(slide_tuples):
        if slide_path is None:
            continue
        if show_progress and idx % 20 == 0:
            print(f"  Slide {idx+1}/{total}...", file=sys.stderr, flush=True)
        try:
            top_text = ocr_region(slide_path, top_frac=0.0, bottom_frac=0.22)
            if is_answer_slide(top_text):
                answer_indices.add(idx)
        except Exception:
            pass

    if show_progress:
        print(f"Found {len(answer_indices)} answer slides.", file=sys.stderr, flush=True)

    # ── Step 3: region OCR on answer slides only ─────────────────────────────
    answer_texts   = {}
    question_texts = {}

    q_indices = set()
    for ans_idx in sorted(answer_indices):
        q_idx = ans_idx - 1
        if q_idx >= 0 and q_idx not in answer_indices:
            q_indices.add(q_idx)

    total_ocr = len(answer_indices) + len(q_indices)
    done = 0

    for ans_idx in sorted(answer_indices):
        slide_path = slide_tuples[ans_idx][0]
        if slide_path is None:
            answer_texts[ans_idx] = ""
            continue
        try:
            answer_texts[ans_idx] = ocr_answer_slide(slide_path)
        except Exception:
            answer_texts[ans_idx] = ""
        done += 1
        if show_progress:
            print(f"  OCR {done}/{total_ocr} (answer slide {ans_idx+1})...", file=sys.stderr, flush=True)

    # ── Step 4: region OCR on question slides only ───────────────────────────
    for q_idx in sorted(q_indices):
        slide_path = slide_tuples[q_idx][0]
        if slide_path is None:
            question_texts[q_idx] = ""
            continue
        try:
            question_texts[q_idx] = ocr_question_slide(slide_path)
        except Exception:
            question_texts[q_idx] = ""
        done += 1
        if show_progress:
            print(f"  OCR {done}/{total_ocr} (question slide {q_idx+1})...", file=sys.stderr, flush=True)

    # ── Step 6: pair Q/A slides and build question list ──────────────────────
    questions = []
    for ans_idx in sorted(answer_indices):
        q_idx = ans_idx - 1
        if q_idx < 0:
            continue

        q_text = question_texts.get(q_idx, "")
        a_text = answer_texts.get(ans_idx, "")

        if not a_text:
            continue  # No answer = nothing useful to save

        # The embedded image from the answer slide is the visual answer
        ans_embed_path = slide_tuples[ans_idx][1]

        # "SOMEONE FAMOUS" / visual-ID slides: question slide just says
        # "answer on next slide" but the ANSWER slide has the person's name.
        # Keep these — use generic question text so answer + image are saved.
        is_visual = is_visual_question(q_text)

        if not q_text or (is_junk_slide(q_text) and not is_visual):
            if show_progress:
                print(f"  Skipping junk slide {q_idx+1}: {q_text[:60]}", file=sys.stderr, flush=True)
            continue

        if is_visual:
            q_text = "Identify this famous personality."

        topic = detect_topic(q_text, a_text, topic_keywords)
        questions.append({
            "question_text":       q_text,
            "answer_text":         a_text,
            "funda_text":          "",
            "topic":               topic,
            "has_image_placeholder": ans_embed_path is not None,
            "question_image_path": None,
            "answer_image_path":   ans_embed_path,  # image on the answer slide = the famous person
        })

    if show_progress:
        print(f"Done. {len(questions)} questions extracted.", file=sys.stderr)

    # temp_dir NOT deleted — questions_updater will copy images from it
    return questions


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ocr_extractor.py <pdf_or_pptx>", file=sys.stderr, flush=True)
        sys.exit(1)

    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
    ROOT_DIR   = os.path.dirname(SCRIPT_DIR)
    with open(os.path.join(ROOT_DIR, "config.json")) as f:
        cfg = json.load(f)

    qs = extract_from_image_file(sys.argv[1], cfg)
    print(json.dumps({
        "title":     os.path.splitext(os.path.basename(sys.argv[1]))[0],
        "questions": qs,
    }, indent=2, ensure_ascii=False))
