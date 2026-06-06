"""
assign_visual_images_v2.py — Smarter visual-image matcher (v2).

Matches broken-visual questions to source PPTX slides using:
- ANSWER text (has distinctive proper nouns)
- Full-slide OCR
- TF-IDF style rarity scoring
- Builds a preview HTML for human verification before auto-applying.
"""
import os, re, sys, json, math, shutil, datetime, subprocess
from collections import Counter

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR   = os.path.dirname(SCRIPT_DIR)
SITE       = r"D:\QUIZBANK"
PROCESSED  = os.path.join(ROOT_DIR, "processed")
IMAGES     = os.path.join(SITE, "images")
PREVIEW    = os.path.join(ROOT_DIR, "visual_match_preview")
RENDERS    = os.path.join(ROOT_DIR, "_slide_renders")
OCR_CACHE  = os.path.join(ROOT_DIR, "_slide_ocr_cache.json")

os.makedirs(PREVIEW, exist_ok=True)
os.makedirs(RENDERS, exist_ok=True)

AUTO_MIN   = 0.40
REVIEW_MIN = 0.22

STOPWORDS = set(("a an and are as at be by for from has have he her him his i in is it its "
                 "me my of on or our she that the their them they this to was we were what "
                 "when where which who will with you your would also more than then over after "
                 "into one two three some all any not but no how why such between under above "
                 "these those each both either neither image picture pictured below shown").split())


def tokenize(text):
    if not text: return []
    text = re.sub(r"<[^>]+>", " ", text).lower()
    return [t for t in re.findall(r"[a-z][a-z']{2,}", text)
            if t not in STOPWORDS and len(t) >= 3]


def load_broken():
    with open(os.path.join(SCRIPT_DIR, "_broken_visuals.json"), encoding="utf-8") as f:
        return json.load(f)


def render_slides(pptx_path):
    """Render each slide to PNG. Returns list of (slide_index, png_path)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io as _io

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    prs = Presentation(pptx_path)
    out = []
    for i, slide in enumerate(prs.slides):
        dest = os.path.join(RENDERS, f"{base}_slide_{i+1:03d}.png")
        if os.path.isfile(dest):
            out.append((i+1, dest))
            continue
        pics = [s for s in slide.sh