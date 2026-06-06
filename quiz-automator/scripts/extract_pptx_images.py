"""
extract_pptx_images.py — Dump every embedded image from every PPTX in processed/,
organized by slide index so we can manually assign them to questions.

Output: quiz-automator/extracted_pptx_images/<pptx_basename>/slide_<NNN>_<MM>.<ext>

Use this when the OCR pipeline failed to extract an image for a question that
clearly needs one (e.g. "Connect the two visuals", "Pictured below is...").
"""
import os
import sys
import io

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR   = os.path.dirname(SCRIPT_DIR)
PROCESSED  = os.path.join(ROOT_DIR, "processed")
OUT_DIR    = os.path.join(ROOT_DIR, "extracted_pptx_images")


def main():
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    os.makedirs(OUT_DIR, exist_ok=True)

    pptx_files = [f for f in os.listdir(PROCESSED) if f.lower().endswith(".pptx")]
    if not pptx_files:
        print(f"No PPTX files found in {PROCESSED}")
        sys.exit(1)

    print(f"Found {len(pptx_files)} PPTX file(s)\n")
    total_imgs = 0

    for fname in pptx_files:
        base = os.path.splitext(fname)[0]
        target_dir = os.path.join(OUT_DIR, base)
        os.makedirs(target_dir, exist_ok=True)

        print(f"=== {fname} ===")
        prs = Presentation(os.path.join(PROCESSED, fname))
        slide_count = 0
        img_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_count += 1
            pic_idx = 0
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                pic_idx += 1
                blob = shape.image.blob
                ext  = shape.image.ext or "png"
                out  = os.path.join(target_dir,
                                    f"slide_{slide_idx+1:03d}_{pic_idx:02d}.{ext}")
                with open(out, "wb") as f:
                    f.write(blob)
                img_count += 1

        print(f"  Slides: {slide_count}, Images extracted: {img_count}\n")
        total_imgs += img_count

    print(f"Total images dumped: {total_imgs}")
    print(f"Output folder:       {OUT_DIR}")


if __name__ == "__main__":
    main()
