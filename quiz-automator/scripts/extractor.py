"""
extractor.py — Extract quiz questions (and embedded images) from PDF, PPTX, ODF, ODT, ODP, or text files.
Usage: python scripts/extractor.py <file_path>
Output: JSON to stdout. Questions may include "question_image_path" pointing to a temp file.
"""
import sys
import os
import re
import json
import shutil
import tempfile
import unicodedata
import zipfile
import xml.etree.ElementTree as ET

SCRIPT_DIR  = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR    = os.path.dirname(SCRIPT_DIR)
CONFIG_PATH = os.path.join(ROOT_DIR, "config.json")

# ODF XML namespaces
ODF_NS = {
    "office":       "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "draw":         "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
    "text":         "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "presentation": "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0",
    "xlink":        "http://www.w3.org/1999/xlink",
    "svg":          "urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0",
}
# Clark notation helpers
def _n(prefix, local):
    return f"{{{ODF_NS[prefix]}}}{local}"


def load_config():
    with open(CONFIG_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


# ── normalisation ────────────────────────────────────────────────────────────
def clean(text):
    text = "".join(c for c in text if unicodedata.category(c) != "Cf")
    text = unicodedata.normalize("NFKC", text)
    return " ".join(text.split())


def is_section_header(line):
    """All-caps line with no colon — e.g. 'SCIENCE QUIZ'. Skipped during parsing."""
    stripped = line.strip()
    if not stripped or ":" in stripped:
        return False
    letters = [c for c in stripped if c.isalpha()]
    if not letters:
        return False
    return all(c.isupper() for c in letters) and not re.match(r"^\d+[\.\)]\s+", stripped)


# ── state machine ────────────────────────────────────────────────────────────
QUESTION_RE = re.compile(r"^\d+[\.\)]\s+(.+)", re.DOTALL)
ANSWER_RE   = re.compile(r"^answer\s*:", re.IGNORECASE)
FUNDA_RE    = re.compile(r"^funda\s*:", re.IGNORECASE)


def parse_questions(lines, image_slots=None):
    """
    Parse lines into question dicts using a state machine.
    image_slots: optional list of (line_index, image_path) — images are assigned
                 to the question that was most recently started before that line index.
    Returns list of question dicts.
    """
    questions = []
    state = "IDLE"
    q_text, a_text, f_text = [], [], []
    question_start_lines = []  # line index where each question started
    current_line = 0

    def flush():
        if q_text:
            questions.append({
                "question_text": clean(" ".join(q_text)),
                "answer_text":   clean(" ".join(a_text)),
                "funda_text":    clean(" ".join(f_text)),
                "question_image_path": None,
            })

    for raw in lines:
        line = clean(raw)
        current_line += 1
        if not line:
            continue
        if is_section_header(line):
            continue

        m = QUESTION_RE.match(line)
        if m:
            flush()
            question_start_lines.append(current_line)
            q_text = [m.group(1).strip()]
            a_text = []
            f_text = []
            state = "READING_QUESTION"
            continue

        if ANSWER_RE.match(line):
            after = ANSWER_RE.sub("", line).strip()
            a_text = [after] if after else []
            state = "READING_ANSWER"
            continue

        if FUNDA_RE.match(line):
            after = FUNDA_RE.sub("", line).strip()
            f_text = [after] if after else []
            state = "READING_FUNDA"
            continue

        if state == "READING_QUESTION":
            q_text.append(line)
        elif state == "READING_ANSWER":
            a_text.append(line)
        elif state == "READING_FUNDA":
            f_text.append(line)

    flush()

    # Assign images: each image goes to the question whose start line is closest before it
    if image_slots:
        for img_line_idx, img_path in image_slots:
            # Find last question started before this image's line
            target_q = None
            for qi, start in enumerate(question_start_lines):
                if start <= img_line_idx:
                    target_q = qi
            if target_q is not None and target_q < len(questions):
                if questions[target_q]["question_image_path"] is None:
                    questions[target_q]["question_image_path"] = img_path

    return questions


# ── PDF extraction ───────────────────────────────────────────────────────────
def extract_text_from_pdf(path):
    lines = []
    try:
        import fitz
        doc = fitz.open(path)
        for page in doc:
            lines.extend(page.get_text("text").splitlines())
        doc.close()
        if lines:
            return lines
    except Exception:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                lines.extend((page.extract_text() or "").splitlines())
        if lines:
            return lines
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        for page in PdfReader(path).pages:
            lines.extend((page.extract_text() or "").splitlines())
        if lines:
            return lines
    except Exception:
        pass
    try:
        from pdfminer.high_level import extract_text as pm_extract
        lines.extend(pm_extract(path).splitlines())
        if lines:
            return lines
    except Exception:
        pass
    raise RuntimeError(f"All PDF extractors failed for: {path}")


# ── PPTX extraction ──────────────────────────────────────────────────────────
def extract_text_from_pptx(path):
    from pptx import Presentation
    lines = []
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
    return lines


# ── ODF extraction ───────────────────────────────────────────────────────────
def _elem_text_lines(elem):
    """Recursively yield non-empty text strings from an XML element."""
    if elem.text and elem.text.strip():
        yield elem.text.strip()
    for child in elem:
        yield from _elem_text_lines(child)
        if child.tail and child.tail.strip():
            yield child.tail.strip()


def _extract_images_from_zip(zf, temp_dir):
    """
    Extract all image files from the ODF zip's Pictures/ folder.
    Returns dict: {internal_href -> local_temp_path}
    e.g. {"Pictures/image1.png": "/tmp/xxx/image1.png"}
    """
    image_map = {}
    valid_ext = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg", ".tiff"}
    for name in zf.namelist():
        if name.lower().startswith("pictures/"):
            ext = os.path.splitext(name)[1].lower()
            if ext in valid_ext:
                dest = os.path.join(temp_dir, os.path.basename(name))
                # Handle duplicate filenames
                if os.path.exists(dest):
                    base, e = os.path.splitext(os.path.basename(name))
                    dest = os.path.join(temp_dir, f"{base}_{len(image_map)}{e}")
                with zf.open(name) as src, open(dest, "wb") as dst:
                    dst.write(src.read())
                image_map[name] = dest
    return image_map


def extract_from_odp(path, temp_dir):
    """
    Extract from ODP (presentation). Each slide is treated as a unit.
    Returns (lines, image_slots) where image_slots = [(line_index, path), ...]
    """
    lines = []
    image_slots = []

    with zipfile.ZipFile(path, "r") as zf:
        image_map = _extract_images_from_zip(zf, temp_dir)
        content_xml = zf.read("content.xml")

    root = ET.fromstring(content_xml)

    # Find all draw:page elements (slides)
    pages = root.findall(f".//{_n('draw','page')}")

    for page in pages:
        slide_start_line = len(lines)

        # Collect all text frames on this slide (skip presentation:notes)
        for frame in page.findall(f".//{_n('draw','text-box')}"):
            # Skip notes frames
            parent_frame = None
            for f in page.iter():
                if f.tag == _n("draw", "text-box") and f is frame:
                    break
            for text_p in frame.findall(f".//{_n('text','p')}"):
                text = " ".join(_elem_text_lines(text_p))
                if clean(text):
                    lines.append(text)

        # Find images on this slide
        for img_elem in page.findall(f".//{_n('draw','image')}"):
            href = img_elem.get(_n("xlink", "href"), "")
            if href in image_map:
                image_slots.append((len(lines), image_map[href]))
            elif href:
                # Try matching by basename
                basename = os.path.basename(href)
                for key, path_val in image_map.items():
                    if os.path.basename(key) == basename:
                        image_slots.append((len(lines), path_val))
                        break

    return lines, image_slots


def extract_from_odt(path, temp_dir):
    """
    Extract from ODT (text document). Images are inline with text.
    Returns (lines, image_slots) where image_slots = [(line_index, path), ...]
    """
    lines = []
    image_slots = []

    with zipfile.ZipFile(path, "r") as zf:
        image_map = _extract_images_from_zip(zf, temp_dir)
        content_xml = zf.read("content.xml")

    root = ET.fromstring(content_xml)
    body = root.find(f".//{_n('office','text')}")
    if body is None:
        body = root

    def process_element(elem):
        """Walk elements in document order, yielding (type, value).
        type: 'text' | 'image'
        """
        tag = elem.tag
        # Paragraph
        if tag == _n("text", "p") or tag == _n("text", "h"):
            # First check for inline images inside this paragraph
            for child in elem.iter():
                if child.tag == _n("draw", "image"):
                    href = child.get(_n("xlink", "href"), "")
                    if href in image_map:
                        yield ("image", image_map[href])
                    elif href:
                        basename = os.path.basename(href)
                        for key, pv in image_map.items():
                            if os.path.basename(key) == basename:
                                yield ("image", pv)
                                break
            # Then extract the text
            text = " ".join(_elem_text_lines(elem))
            if clean(text):
                yield ("text", text)
        else:
            for child in elem:
                yield from process_element(child)

    for item_type, value in process_element(body):
        if item_type == "text":
            lines.append(value)
        elif item_type == "image":
            image_slots.append((len(lines), value))

    return lines, image_slots


def extract_odf(path, temp_dir):
    """Dispatch to ODP or ODT extractor based on mimetype in the ODF zip."""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            try:
                mimetype = zf.read("mimetype").decode("utf-8").strip()
            except Exception:
                mimetype = ""

        if "presentation" in mimetype or path.lower().endswith(".odp"):
            return extract_from_odp(path, temp_dir)
        else:
            # Default to ODT (text document) for .odt, .odf, or unknown
            return extract_from_odt(path, temp_dir)
    except zipfile.BadZipFile:
        # Not a zip — try as plain text
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read().splitlines(), []


# ── topic detection ──────────────────────────────────────────────────────────
def run_image_extractor(path, config):
    """
    Pick the best available image-extraction backend:
    - Gemini Vision API (if gemini_api_key in config) — high quality
    - EasyOCR fallback                                  — free, local, lower quality
    """
    api_key = (config.get("gemini_api_key") or "").strip()
    if api_key:
        print("No text found — using Gemini Vision API for extraction...", file=sys.stderr, flush=True)
        try:
            from gemini_extractor import extract_from_image_file as gemini_extract
            return gemini_extract(path, config)
        except Exception as e:
            print(f"Gemini extraction failed ({e}); falling back to local OCR...", file=sys.stderr, flush=True)

    print("Using local EasyOCR for extraction...", file=sys.stderr, flush=True)
    from ocr_extractor import extract_from_image_file as ocr_extract
    return ocr_extract(path, config)


def detect_topic(question_text, answer_text, funda_text, topic_keywords):
    combined = " ".join([question_text, answer_text, funda_text]).lower()
    scores = {
        topic: sum(1 for kw in kws if kw.lower() in combined)
        for topic, kws in topic_keywords.items()
    }
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "general"


# ── main ─────────────────────────────────────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print("Usage: python scripts/extractor.py <file_path>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.isfile(path):
        print(f"File not found: {path}", file=sys.stderr)
        sys.exit(1)

    try:
        config = load_config()
    except Exception as e:
        print(f"Failed to load config.json: {e}", file=sys.stderr)
        sys.exit(1)

    topic_keywords = config.get("topic_keywords", {})
    title = os.path.splitext(os.path.basename(path))[0]
    ext = os.path.splitext(path)[1].lower()

    temp_dir = tempfile.mkdtemp(prefix="quiz_extract_")
    try:
        raw_questions = None

        if ext in (".odf", ".odp", ".odt"):
            lines, image_slots = extract_odf(path, temp_dir)
            raw_questions = parse_questions(lines, image_slots)

        elif ext == ".pdf":
            lines = extract_text_from_pdf(path)
            meaningful = [l for l in lines if len(l.strip()) > 3]
            if len(meaningful) < 5:
                raw_questions = run_image_extractor(path, config)
            else:
                raw_questions = parse_questions(lines, [])

        elif ext in (".pptx", ".ppt"):
            lines = extract_text_from_pptx(path)
            meaningful = [l for l in lines if len(l.strip()) > 3]
            if len(meaningful) < 5:
                raw_questions = run_image_extractor(path, config)
            else:
                raw_questions = parse_questions(lines, [])

        else:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                lines = f.read().splitlines()
            raw_questions = parse_questions(lines, [])

        if not raw_questions:
            print("No questions found in file.", file=sys.stderr)
            sys.exit(1)

        output_questions = []
        for q in raw_questions:
            # OCR extractor already sets topic; text extractor needs it computed
            topic = q.get("topic") or detect_topic(
                q["question_text"], q["answer_text"], q.get("funda_text", ""), topic_keywords
            )
            img_path = q.get("question_image_path")
            output_questions.append({
                "question_text":         q["question_text"],
                "answer_text":           q["answer_text"],
                "funda_text":            q.get("funda_text", ""),
                "topic":                 topic,
                "has_image_placeholder": img_path is not None,
                "question_image_path":   img_path,
            })

        result = {"title": title, "questions": output_questions}
        print(json.dumps(result, ensure_ascii=False, indent=2))

    except Exception as e:
        # Clean up temp dir on failure
        shutil.rmtree(temp_dir, ignore_errors=True)
        print(f"Extraction error: {e}", file=sys.stderr)
        sys.exit(1)
    # Note: temp_dir is NOT cleaned up here — questions_updater needs to copy images from it.
    # watcher.py cleans it up after questions_updater finishes.


if __name__ == "__main__":
    main()
